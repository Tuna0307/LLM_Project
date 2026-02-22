"""
ingest.py - ETL Pipeline for processing learning materials.

Handles:
- Multi-format text extraction: PDF, DOCX, PPTX, TXT, MD, HTML, CSV
- Regex-based noise removal (headers, footers, watermarks)
- Structure-aware semantic chunking with heading detection
- Metadata enrichment (source file, page number, topic, week, doc type)

Author: Wei Xuan (Data Infrastructure)
"""
import os
import re
from typing import Optional

from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document

import config

# All file extensions the pipeline can handle
SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".pptx",
    ".txt", ".md", ".markdown",
    ".html", ".htm", ".csv",
}


# ── Noise Patterns to Remove ────────────────────────────────────────────────
NOISE_PATTERNS = [
    r"©\s*\d{4}.*?reserved\.?",                  # Copyright notices
    r"Singapore Institute of Technology",          # SIT watermarks
    r"All [Rr]ights [Rr]eserved",                 # Rights reserved
    r"Page\s*\d+\s*(of\s*\d+)?",                  # Page numbers
    r"^\s*\d+\s*$",                                # Standalone page numbers
    r"Confidential|CONFIDENTIAL",                  # Confidentiality notices
    r"^\s*[-–—]+\s*$",                             # Separator lines
]

# Compiled regex for efficiency
_noise_re = re.compile("|".join(NOISE_PATTERNS), re.IGNORECASE | re.MULTILINE)

# Heading detection for structure-aware chunking
_heading_re = re.compile(config.HEADING_PATTERN, re.IGNORECASE | re.MULTILINE)


def clean_text(text: str) -> str:
    """Remove noise from extracted PDF text."""
    # Remove noise patterns
    cleaned = _noise_re.sub("", text)
    # Collapse multiple blank lines into one
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    # Strip leading/trailing whitespace
    return cleaned.strip()


def detect_heading(text: str) -> Optional[str]:
    """Detect section headings in text for structure-aware chunking."""
    match = _heading_re.search(text)
    if match:
        # Extract the full line containing the heading
        start = text.rfind("\n", 0, match.start()) + 1
        end = text.find("\n", match.end())
        if end == -1:
            end = len(text)
        return text[start:end].strip()
    return None


def extract_metadata_from_filename(filename: str) -> dict:
    """
    Extract metadata hints from the filename.
    Expected patterns: 'Week3_Preprocessing.pdf', 'Lecture_5_NLP.pdf', etc.
    """
    metadata = {"doc_type": "lecture"}  # default

    # Try to extract week number
    week_match = re.search(r"[Ww]eek\s*(\d+)", filename)
    if week_match:
        metadata["week"] = int(week_match.group(1))

    # Try to extract topic from filename
    name_no_ext = os.path.splitext(filename)[0]
    # Remove week/lecture prefixes to get topic
    topic = re.sub(r"[Ww]eek\s*\d+[_\s-]*", "", name_no_ext)
    topic = re.sub(r"[Ll]ecture\s*\d+[_\s-]*", "", topic)
    topic = topic.replace("_", " ").replace("-", " ").strip()
    if topic:
        metadata["topic"] = topic

    # Detect document type
    lower_name = filename.lower()
    if "tutorial" in lower_name or "tut" in lower_name:
        metadata["doc_type"] = "tutorial"
    elif "lab" in lower_name:
        metadata["doc_type"] = "lab"
    elif "handbook" in lower_name or "guide" in lower_name:
        metadata["doc_type"] = "handbook"

    return metadata


# ── Per-format Page Extractors ─────────────────────────────────────────────
# Each returns a list of (text, page_number) tuples.

def _extract_pages_pdf(file_path: str) -> list[tuple[str, int]]:
    """Extract text page-by-page from a PDF."""
    loader = PyPDFLoader(file_path)
    raw_pages = loader.load()
    return [
        (page.page_content, page.metadata.get("page", i) + 1)
        for i, page in enumerate(raw_pages)
    ]


def _extract_pages_docx(file_path: str) -> list[tuple[str, int]]:
    """Extract text from a Word (.docx) file, simulating pages."""
    import docx as _docx
    doc = _docx.Document(file_path)
    texts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also pull text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                texts.append(row_text)
    # Group into simulated pages (every 30 paragraphs)
    PAGE_SIZE = 30
    pages: list[tuple[str, int]] = []
    for i in range(0, max(len(texts), 1), PAGE_SIZE):
        page_text = "\n".join(texts[i:i + PAGE_SIZE])
        if page_text.strip():
            pages.append((page_text, i // PAGE_SIZE + 1))
    return pages


def _extract_pages_pptx(file_path: str) -> list[tuple[str, int]]:
    """Extract text slide-by-slide from a PowerPoint (.pptx) file."""
    from pptx import Presentation
    prs = Presentation(file_path)
    pages: list[tuple[str, int]] = []
    for i, slide in enumerate(prs.slides):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        # Include speaker notes if present
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Notes] {notes}")
        if parts:
            pages.append(("\n".join(parts), i + 1))
    return pages


def _extract_pages_text(file_path: str) -> list[tuple[str, int]]:
    """Extract text from TXT, MD, CSV, HTML files."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()

    ext = os.path.splitext(file_path)[1].lower()
    if ext in (".html", ".htm"):
        # Strip HTML tags using the stdlib parser
        from html.parser import HTMLParser
        class _Extractor(HTMLParser):
            def __init__(self) -> None:
                super().__init__()
                self.parts: list[str] = []
                self._skip = False
            def handle_starttag(self, tag, attrs):
                if tag in ("script", "style"):
                    self._skip = True
            def handle_endtag(self, tag):
                if tag in ("script", "style"):
                    self._skip = False
            def handle_data(self, data):
                if not self._skip:
                    self.parts.append(data)
        parser = _Extractor()
        parser.feed(content)
        content = " ".join(parser.parts)

    # Group into simulated pages of ~3000 characters
    PAGE_SIZE = 3000
    pages: list[tuple[str, int]] = []
    for i in range(0, max(len(content), 1), PAGE_SIZE):
        chunk = content[i:i + PAGE_SIZE]
        if chunk.strip():
            pages.append((chunk, i // PAGE_SIZE + 1))
    return pages


# ── Unified File Processor ───────────────────────────────────────────────────

def load_and_process_file(
    file_path: str,
    extra_metadata: Optional[dict] = None,
    original_filename: Optional[str] = None,
) -> list[Document]:
    """
    Full ETL pipeline for any supported file type.

    Supported: .pdf, .docx, .pptx, .txt, .md, .markdown, .html, .htm, .csv

    1. Extract text (format-aware)
    2. Clean noise
    3. Split into semantic chunks with overlap
    4. Enrich with metadata

    Args:
        file_path: Path to the file
        extra_metadata: Optional additional metadata (e.g., user-provided topic)
        original_filename: The real filename to store in metadata (use when
            file_path points to a temp file with a generated name)

    Returns:
        List of Document objects ready for embedding
    """
    ext = os.path.splitext(file_path)[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type '{ext}'. "
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    # Step 1: Extract raw pages with the right loader
    if ext == ".pdf":
        raw_pages = _extract_pages_pdf(file_path)
    elif ext == ".docx":
        raw_pages = _extract_pages_docx(file_path)
    elif ext == ".pptx":
        raw_pages = _extract_pages_pptx(file_path)
    else:  # .txt, .md, .markdown, .html, .htm, .csv
        raw_pages = _extract_pages_text(file_path)

    filename = original_filename if original_filename else os.path.basename(file_path)
    file_metadata = extract_metadata_from_filename(filename)
    if extra_metadata:
        file_metadata.update(extra_metadata)

    # Step 2 & 3: Clean + chunk
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=config.CHUNK_SIZE,
        chunk_overlap=config.CHUNK_OVERLAP,
        separators=["\n\n", "\n", ". ", " ", ""],
        length_function=len,
    )

    all_chunks: list[Document] = []
    current_heading = "Introduction"

    for page_text, page_num in raw_pages:
        cleaned_text = clean_text(page_text)
        if not cleaned_text:
            continue

        heading = detect_heading(cleaned_text)
        if heading:
            current_heading = heading

        page_chunks = text_splitter.split_text(cleaned_text)
        for i, chunk_text in enumerate(page_chunks):
            contextualized_text = f"[{current_heading}]\n{chunk_text}"
            chunk_metadata = {
                "source_file": filename,
                "page_number": page_num,
                "chunk_index": i,
                "section_heading": current_heading,
                **file_metadata,
            }
            all_chunks.append(
                Document(page_content=contextualized_text, metadata=chunk_metadata)
            )

    return all_chunks


# Backward compatibility alias
load_and_process_pdf = load_and_process_file


def ingest_directory(
    directory: Optional[str] = None,
    extra_metadata: Optional[dict] = None,
) -> list[Document]:
    """
    Process all PDF files in a directory.

    Args:
        directory: Path to directory containing PDFs. Defaults to config.RAW_DATA_DIR.
        extra_metadata: Optional metadata to apply to all files.

    Returns:
        List of all Document chunks from all PDFs.
    """
    if directory is None:
        directory = config.RAW_DATA_DIR

    all_documents = []
    all_files = [
        f for f in os.listdir(directory)
        if os.path.splitext(f)[1].lower() in SUPPORTED_EXTENSIONS
    ]

    if not all_files:
        print(f"[WARN] No supported files found in {directory}")
        return all_documents

    for filename in sorted(all_files):
        file_path = os.path.join(directory, filename)
        print(f"[>>] Processing: {filename}")

        try:
            docs = load_and_process_file(file_path, extra_metadata)
            all_documents.extend(docs)
            print(f"   [OK] Generated {len(docs)} chunks")
        except Exception as e:
            print(f"   [ERR] Error processing {filename}: {e}")

    print(f"\n[OK] Total: {len(all_documents)} chunks from {len(all_files)} files")
    return all_documents


if __name__ == "__main__":
    # CLI usage: python -m src.ingest
    docs = ingest_directory()
    if docs:
        print(f"\nSample chunk metadata: {docs[0].metadata}")
        print(f"Sample chunk text (first 200 chars): {docs[0].page_content[:200]}")
